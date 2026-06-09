#!/usr/bin/env python3
"""Generate a 10-slide Agent Skill Router PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette ───────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_DARK = RGBColor(0x0F, 0x17, 0x2A)
GRAY = RGBColor(0x47, 0x55, 0x69)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED_ACCENT = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
TEAL = RGBColor(0x14, 0xB8, 0xA6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FOOTER_H = Inches(0.5)


def px(size):
    return Pt(size)


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
        else:
            shape.line.width = Pt(1)
    return shape


def add_rounded_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
        else:
            shape.line.width = Pt(1)
    return shape


def add_circle(slide, x, y, size, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
        else:
            shape.line.width = Pt(1)
    return shape


def add_textbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(x, y, w, h)


def set_text(shape_or_tb, text, size=14, bold=False, color=SLATE_DARK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    if hasattr(shape_or_tb, "text_frame"):
        tf = shape_or_tb.text_frame
    else:
        tf = shape_or_tb
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size) if isinstance(size, int) else size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return p


def add_para(tf, text, size=14, bold=False, color=SLATE_DARK, align=PP_ALIGN.LEFT, space_before=Pt(4), font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size) if isinstance(size, int) else size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    if space_before:
        p.space_before = space_before
    return p


def add_icon_shape(slide, x, y, size, fill, icon_text):
    """Add a colored circle with white text as an icon."""
    circ = add_circle(slide, x, y, size, fill=fill)
    tb = add_textbox(slide, x, y, size, size)
    set_text(tb, icon_text, size=int(size/Inches(1)*14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return circ


def add_footer(slide, text="Agent Skill Router"):
    bar = add_rect(slide, Inches(0), Inches(7.5) - FOOTER_H, SLIDE_W, FOOTER_H, fill=NAVY)
    tb = add_textbox(slide, Inches(0.4), Inches(7.5) - FOOTER_H + Pt(3), Inches(8), Inches(FOOTER_H))
    set_text(tb, text, size=9, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT)


def add_section_header(slide, title):
    """Add a navy bar with white title at top."""
    bar = add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), fill=NAVY)
    tb = add_textbox(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.9))
    set_text(tb, title, size=36, bold=True, color=WHITE, font_name="Segoe UI")


def bullet_list(slide, x, y, items, sizes=None):
    tb = add_textbox(slide, x, y, Inches(6), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        sz = sizes[i] if sizes else Pt(18)
        p.text = "▸ " + item
        p.font.size = sz
        p.font.color.rgb = SLATE_DARK
        p.font.name = "Calibri"
        p.space_before = Pt(8)
    return tb


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════
def slide_title(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Full navy background
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=NAVY)

    # Decorative accent shapes
    add_rounded_rect(slide, Inches(-1), Inches(-1), Inches(4), Inches(4), fill=BLUE)
    add_circle(slide, Inches(10.5), Inches(5.5), Inches(3), fill=RGBColor(0x1E, 0x3A, 0x5F))
    add_rounded_rect(slide, Inches(9), Inches(-1), Inches(2.5), Inches(2.5), fill=BLUE)

    # AI Brain icon circle in center-left area
    circ = add_circle(slide, Inches(7.5), Inches(0.8), Inches(2.5), fill=BLUE)
    tb_icon = add_textbox(slide, Inches(7.5), Inches(0.8), Inches(2.5), Inches(2.5))
    set_text(tb_icon, "⚡", size=64, align=PP_ALIGN.CENTER)

    # Title
    tb = add_textbox(slide, Inches(1), Inches(2.0), Inches(10), Inches(1.5))
    set_text(tb, "Agent Skill Router", size=52, bold=True, color=WHITE, font_name="Segoe UI", align=PP_ALIGN.LEFT)

    # Subtitle
    tb = add_textbox(slide, Inches(1), Inches(3.6), Inches(9), Inches(0.7))
    set_text(tb, "Intelligent Skill Routing for AI Agents", size=26, color=BLUE_LIGHT, font_name="Segoe UI", align=PP_ALIGN.LEFT)

    # One-liner
    tb = add_textbox(slide, Inches(1), Inches(4.5), Inches(9), Inches(0.6))
    set_text(tb, 'Just-in-time AI expertise. Routed automatically.', size=18, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT)

    # Bottom accent line
    add_rect(slide, Inches(1), Inches(5.5), Inches(2.5), Pt(3), fill=BLUE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS IT?
# ═══════════════════════════════════════════════════════════════════
def slide_what_is(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_section_header(slide, "What Is the Agent Skill Router?")
    add_footer(slide)

    # Left visual: Big stat in colored circle
    big_circ = add_circle(slide, Inches(0.8), Inches(1.8), Inches(3.5), fill=BLUE)
    tb_num = add_textbox(slide, Inches(0.8), Inches(2.2), Inches(3.5), Inches(1.2))
    set_text(tb_num, "1,236", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Segoe UI")

    tb_label = add_textbox(slide, Inches(0.8), Inches(3.8), Inches(3.5), Inches(0.5))
    set_text(tb_label, "Skills Loaded", size=16, color=WHITE, align=PP_ALIGN.CENTER)

    # Decorative ring around number
    ring = add_circle(slide, Inches(0.8), Inches(1.8), Inches(3.5), fill=None, line=BLUE_LIGHT, line_w=Pt(2))
    ring.fill.background()

    # Right text: bullets
    items = [
        "Embeds tasks into vector space",
        "Retrieves matching skills via hybrid scoring",
        "Delivers expert knowledge to AI agents",
    ]
    bullet_list(slide, Inches(5.2), Inches(1.8), items)

    # Small label at top-right of text area
    tb_badge = add_textbox(slide, Inches(5.2), Inches(1.4), Inches(6), Inches(0.4))
    set_text(tb_badge, "How it works — in brief", size=13, color=GRAY, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — THE REQUEST FLOW
# ═══════════════════════════════════════════════════════════════════
def slide_request_flow(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_section_header(slide, "The Request Flow")
    add_footer(slide)

    steps = [
        ("User\nquery", 0),
        ("Task\nembedding", BLUE),
        ("Hybrid\nsearch", GREEN),
        ("LLM\nranking", AMBER),
        ("Skill\nfetch", PURPLE),
        ("Context\ninjection", TEAL),
    ]

    total_w = Inches(11.5)
    step_w = Inches(1.7)
    gap = (total_w - 6 * step_w) / 5
    start_x = Inches(0.8)

    for i, (label, color) in enumerate(steps):
        x = start_x + i * (step_w + gap)

        # Circle
        circ = add_circle(slide, x + Inches(0.3), Inches(1.5), Inches(1.1), fill=NAVY if i == 0 else BLUE if color == 0 else color)
        circ.fill.solid()

        # Number text
        tb = add_textbox(slide, x + Inches(0.3), Inches(1.5), Inches(1.1), Inches(1.1))
        set_text(tb, str(i + 1), size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Arrow between circles (except last)
        if i < 5:
            arrow_x = x + step_w + Pt(4)
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, Inches(1.9), Inches(0.5), Inches(0.3))
            arr.fill.solid()
            arr.fill.fore_color.rgb = BLUE_LIGHT
            arr.line.fill.background()

        # Label below
        tb_lbl = add_textbox(slide, x, Inches(2.8), step_w, Inches(1.5))
        set_text(tb_lbl, label, size=13, bold=True, color=SLATE_DARK, align=PP_ALIGN.CENTER)

    # Bottom performance stat
    perf_bar = add_rounded_rect(slide, Inches(3.5), Inches(4.8), Inches(6.3), Inches(0.8), fill=BLUE_LIGHT, line=BLUE, line_w=Pt(1))
    tb_perf = add_textbox(slide, Inches(3.5), Inches(4.85), Inches(6.3), Inches(0.7))
    set_text(tb_perf, "⚡  ~10ms warm   |   ~3.5s cold", size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — SIX-STAGE PIPELINE
# ═══════════════════════════════════════════════════════════════════
def slide_pipeline(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_section_header(slide, "Six-Stage Pipeline")
    add_footer(slide)

    stages = [
        ("Safety", "~0.1ms", GREEN, "🛡"),
        ("Search", "~200ms", BLUE, "🔍"),
        ("Match", "~0.1ms", AMBER, "✓"),
        ("Diversify", "~0.1ms", PURPLE, "⚖"),
        ("Rank", "~3s", RED_ACCENT, "★"),
        ("Plan", "~0.1ms", TEAL, "📋"),
    ]

    cols = 3
    rows = 2
    card_w = Inches(3.2)
    card_h = Inches(2.0)
    gap_x = Inches(0.5)
    gap_y = Inches(0.4)
    start_x = Inches(1.1)
    start_y = Inches(1.6)

    for idx, (name, time, color, icon) in enumerate(stages):
        col = idx % cols
        row = idx // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        # Card background
        card = add_rounded_rect(slide, x, y, card_w, card_h, fill=LIGHT_BG, line=RGBColor(0xE2, 0xE8, 0xF0), line_w=Pt(1))

        # Color top bar
        add_rect(slide, x, y, card_w, Inches(0.08), fill=color)

        # Icon circle
        ic = add_circle(slide, x + Inches(1.15), y + Inches(0.35), Inches(0.7), fill=color)

        # Name
        tb_name = add_textbox(slide, x + Inches(0.2), y + Inches(1.1), card_w - Inches(0.4), Inches(0.4))
        set_text(tb_name, name, size=18, bold=True, color=SLATE_DARK, align=PP_ALIGN.CENTER)

        # Time
        tb_time = add_textbox(slide, x + Inches(0.2), y + Inches(1.5), card_w - Inches(0.4), Inches(0.3))
        set_text(tb_time, time, size=13, color=GRAY, align=PP_ALIGN.CENTER)

    # Bottom stat callout
    stat_circ = add_circle(slide, Inches(5.5), Inches(5.8), Inches(2.3), fill=BLUE)
    tb_stat_num = add_textbox(slide, Inches(5.5), Inches(6.0), Inches(2.3), Inches(0.6))
    set_text(tb_stat_num, "84%", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Segoe UI")

    tb_stat_lbl = add_textbox(slide, Inches(5.5), Inches(6.45), Inches(2.3), Inches(0.4))
    set_text(tb_stat_lbl, "Cache Hit Rate", size=12, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — SKILLS: THE KNOWLEDGE BASE
# ═══════════════════════════════════════════════════════════════════
def slide_skills_db(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_section_header(slide, "Skills: The Knowledge Base")
    add_footer(slide)

    # Left visual: stacked document shapes
    base_x = Inches(1.0)
    base_y = Inches(2.0)
    doc_w = Inches(3.5)
    doc_h = Inches(1.2)
    stack_colors = [NAVY, BLUE, BLUE_LIGHT]

    for i in range(3):
        y_off = base_y + i * Inches(0.15)
        x_off = Inches(0.4) * (3 - i)
        doc = add_rounded_rect(slide, base_x + x_off, base_y + y_off, doc_w, doc_h, fill=stack_colors[i])
        # File icon lines on each doc
        for j in range(2):
            add_rect(slide, base_x + x_off + Inches(0.3), base_y + y_off + Inches(0.4) + j * Inches(0.25), Inches(1.5), Pt(3), fill=WHITE if i < 2 else NAVY)

    # SKILL.md label on top doc
    tb_doc_lbl = add_textbox(slide, base_x + Inches(0.4), base_y + Inches(0.15), Inches(3.5), Inches(0.5))
    set_text(tb_doc_lbl, "SKILL.md", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Domain badge below stack
    domain_badge = add_rounded_rect(slide, base_x + Inches(0.4), base_y + doc_h + Inches(0.8), Inches(2.5), Inches(0.6), fill=BLUE_LIGHT, line=BLUE, line_w=Pt(1))
    tb_domain = add_textbox(slide, base_x + Inches(0.4), base_y + doc_h + Inches(0.8), Inches(2.5), Inches(0.6))
    set_text(tb_domain, "24 domains · 8 categories", size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Right side: bullets
    items = [
        "1,236 files across 24 domains",
        "YAML frontmatter + structured markdown",
        "Zero tolerance for stubs: 3KB minimum",
    ]
    bullet_list(slide, Inches(5.5), Inches(2.0), items)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — HYBRID SCORING ENGINE
# ═══════════════════════════════════════════════════════════════════
def slide_scoring(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_section_header(slide, "Hybrid Scoring Engine")
    add_footer(slide)

    components = [
        ("Vector Similarity", 50, BLUE),
        ("BM25 Term Match", 30, GREEN),
        ("Trigger Keywords", 15, AMBER),
        ("Archetype Alignment", 10, PURPLE),
    ]

    bar_max_w = Inches(8)
    label_w = Inches(2.8)
    start_y = Inches(2.0)
    bar_h = Inches(0.65)
    gap = Inches(0.35)

    for i, (label, pct, color) in enumerate(components):
        y = start_y + i * (bar_h + gap)

        # Label
        tb_lbl = add_textbox(slide, Inches(1.2), y, label_w, bar_h)
        set_text(tb_lbl, label, size=16, bold=True, color=SLATE_DARK, align=PP_ALIGN.RIGHT)

        # Background track
        add_rect(slide, Inches(4.2), y + Inches(0.1), Inches(7.5), bar_h - Inches(0.2), fill=LIGHT_BG)

        # Filled bar (width proportional to percentage)
        fill_w = int(bar_max_w * pct / 100)
        add_rect(slide, Inches(4.2), y + Inches(0.1), fill_w, bar_h - Inches(0.2), fill=color)

        # Percentage label on the right
        tb_pct = add_textbox(slide, Inches(12.0), y, Inches(1.2), bar_h)
        set_text(tb_pct, f"{pct}%", size=20, bold=True, color=color, align=PP_ALIGN.LEFT)

    # Weight note at bottom
    tb_note = add_textbox(slide, Inches(4.2), Inches(5.6), Inches(7), Inches(0.5))
    set_text(tb_note, "5 independent signals → weighted composite score", size=13, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — DYNAMIC DISCOVERY & SELF-HEALING
# ═══════════════════════════════════════════════════════════════════
def slide_discovery(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_section_header(slide, "Dynamic Discovery & Self-Healing")
    add_footer(slide)

    tiles = [
        ("🔗", "Trigger→\nDomain Index", "Auto-discovered", BLUE),
        ("🤖", "Auto-Skill\nCreation", "Self-generating", GREEN),
        ("🔄", "Periodic\nRefresh", "Every 3,600s", AMBER),
        ("📊", "Token\nTracking", "End-to-end", PURPLE),
    ]

    tile_w = Inches(2.5)
    tile_h = Inches(3.8)
    gap = Inches(0.5)
    start_x = Inches(1.2)
    start_y = Inches(1.6)

    for i, (icon, title, subtitle, color) in enumerate(tiles):
        x = start_x + i * (tile_w + gap)

        # Tile card
        card = add_rounded_rect(slide, x, start_y, tile_w, tile_h, fill=LIGHT_BG, line=color, line_w=Pt(2))

        # Color accent bar at top
        add_rect(slide, x, start_y, tile_w, Inches(0.1), fill=color)

        # Icon area circle
        ic_circ = add_circle(slide, x + Inches(0.65), start_y + Inches(0.4), Inches(1.2), fill=color)

        # Icon text
        tb_icon = add_textbox(slide, x + Inches(0.65), start_y + Inches(0.4), Inches(1.2), Inches(1.2))
        set_text(tb_icon, icon, size=36, align=PP_ALIGN.CENTER)

        # Title
        tb_title = add_textbox(slide, x + Inches(0.15), start_y + Inches(1.9), tile_w - Inches(0.3), Inches(0.9))
        set_text(tb_title, title, size=14, bold=True, color=SLATE_DARK, align=PP_ALIGN.CENTER)

        # Subtitle badge
        sub_badge = add_rounded_rect(slide, x + Inches(0.5), start_y + Inches(3.0), tile_w - Inches(1.0), Inches(0.45), fill=color)
        tb_sub = add_textbox(slide, x + Inches(0.5), start_y + Inches(3.05), tile_w - Inches(1.0), Inches(0.45))
        set_text(tb_sub, subtitle, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — SAFETY & SECURITY
# ═══════════════════════════════════════════════════════════════════
def slide_safety(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_section_header(slide, "Safety & Security")
    add_footer(slide)

    # Left: Shield icon in large colored circle
    shield_circ = add_circle(slide, Inches(1.5), Inches(1.8), Inches(3.0), fill=GREEN)
    tb_shield = add_textbox(slide, Inches(1.5), Inches(1.8), Inches(3.0), Inches(3.0))
    set_text(tb_shield, "🛡", size=80, align=PP_ALIGN.CENTER)

    # Decorative ring
    shield_ring = add_circle(slide, Inches(1.5), Inches(1.8), Inches(3.4), fill=None, line=GREEN, line_w=Pt(2))
    shield_ring.fill.background()

    # Right: bullets
    items = [
        "Two-tier injection detection",
        "Requires 2+ threat signals to block",
    ]
    bullet_list(slide, Inches(5.5), Inches(2.2), items)

    # Label
    tb_label = add_textbox(slide, Inches(5.5), Inches(1.7), Inches(6), Inches(0.4))
    set_text(tb_label, "Protecting AI agents from prompt injection", size=13, color=GRAY)

    # Bottom badge
    badge = add_rounded_rect(slide, Inches(5.5), Inches(4.8), Inches(7), Inches(0.8), fill=RGBColor(0xFE, 0xF3, 0xC7), line=AMBER, line_w=Pt(1))
    tb_badge = add_textbox(slide, Inches(5.5), Inches(4.85), Inches(7), Inches(0.7))
    set_text(tb_badge, '⚠ SAFETY_STRICT=true → single signal triggers block', size=16, bold=True, color=RGBColor(0x92, 0x40, 0x0E))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — TOKEN EFFICIENCY VIA COMPRESSION
# ═══════════════════════════════════════════════════════════════════
def slide_compression(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_section_header(slide, "Token Efficiency via Compression")
    add_footer(slide)

    stats = [
        ("84%", "Cache Hit Rate", BLUE),
        ("65%", "Max Token Savings", GREEN),
        ("3-tier", "Memory → Disk → GitHub", AMBER),
        ("~1.27s", "Startup Time", PURPLE),
    ]

    cols = 2
    cards = []
    card_w = Inches(4.5)
    card_h = Inches(2.0)
    gap_x = Inches(1.0)
    start_x = Inches(1.8)

    for idx, (num, label, color) in enumerate(stats):
        col = idx % cols
        row = idx // cols
        x = start_x + col * (card_w + gap_x)
        y = Inches(2.0) + row * (card_h + Inches(0.5))

        card = add_rounded_rect(slide, x, y, card_w, card_h, fill=LIGHT_BG, line=color, line_w=Pt(2))

        # Number
        tb_num = add_textbox(slide, x + Inches(0.3), y + Inches(0.25), card_w - Inches(0.6), Inches(0.9))
        set_text(tb_num, num, size=40, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Segoe UI")

        # Label
        tb_lbl = add_textbox(slide, x + Inches(0.3), y + Inches(1.2), card_w - Inches(0.6), Inches(0.5))
        set_text(tb_lbl, label, size=14, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — GETTING STARTED
# ═══════════════════════════════════════════════════════════════════
def slide_getting_started(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_section_header(slide, "Get Started in Seconds")
    add_footer(slide)

    steps = [
        ("🚀", "Clone &\nInstall", "./install-skill-router.sh", BLUE),
        ("⚙", "Configure\nAPI Keys", "Interactive prompts", GREEN),
        ("✅", "Auto-route\nStarts", "route_to_skill() enabled", TEAL),
    ]

    card_w = Inches(3.2)
    card_h = Inches(3.5)
    gap = Inches(0.85)
    start_x = Inches(1.5)
    start_y = Inches(1.7)

    for i, (icon, title, detail, color) in enumerate(steps):
        x = start_x + i * (card_w + gap)

        # Card background
        card = add_rounded_rect(slide, x, start_y, card_w, card_h, fill=LIGHT_BG, line=color, line_w=Pt(2))

        # Color top bar
        add_rect(slide, x, start_y, card_w, Inches(0.1), fill=color)

        # Step number circle
        num_circ = add_circle(slide, x + Inches(1.15), start_y + Inches(0.35), Inches(0.65), fill=color)

        # Number text
        tb_num = add_textbox(slide, x + Inches(1.15), start_y + Inches(0.35), Inches(0.65), Inches(0.65))
        set_text(tb_num, str(i + 1), size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Icon
        tb_icon = add_textbox(slide, x + Inches(0.9), start_y + Inches(1.1), Inches(1.5), Inches(0.6))
        set_text(tb_icon, icon, size=32, align=PP_ALIGN.CENTER)

        # Title
        tb_title = add_textbox(slide, x + Inches(0.3), start_y + Inches(1.8), card_w - Inches(0.6), Inches(0.7))
        set_text(tb_title, title, size=16, bold=True, color=SLATE_DARK, align=PP_ALIGN.CENTER)

        # Code/detail box
        detail_box = add_rounded_rect(slide, x + Inches(0.3), start_y + Inches(2.7), card_w - Inches(0.6), Inches(0.55), fill=color if i == 0 else BLUE_LIGHT, line=color, line_w=Pt(1))
        tb_detail = add_textbox(slide, x + Inches(0.3), start_y + Inches(2.75), card_w - Inches(0.6), Inches(0.5))
        set_text(tb_detail, detail, size=12, color=color if i == 0 else NAVY, align=PP_ALIGN.CENTER)

    # Bottom bar
    bottom_bar = add_rounded_rect(slide, Inches(2.5), Inches(6.0), Inches(8.3), Inches(0.6), fill=NAVY)
    tb_bottom = add_textbox(slide, Inches(2.5), Inches(6.05), Inches(8.3), Inches(0.5))
    set_text(tb_bottom, "API: localhost:3000  |  Docker + systemd ready", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# MAIN — BUILD THE PRESENTATION
# ═══════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Build all slides
    slide_title(prs)        # Slide 1 — Title
    slide_what_is(prs)      # Slide 2 — What Is It?
    slide_request_flow(prs) # Slide 3 — Request Flow
    slide_pipeline(prs)     # Slide 4 — Six-Stage Pipeline
    slide_skills_db(prs)    # Slide 5 — Knowledge Base
    slide_scoring(prs)      # Slide 6 — Hybrid Scoring
    slide_discovery(prs)    # Slide 7 — Discovery & Self-Healing
    slide_safety(prs)       # Slide 8 — Safety & Security
    slide_compression(prs)  # Slide 9 — Token Efficiency
    slide_getting_started(prs)  # Slide 10 — Getting Started

    out = "/home/paulpas/git/agent-skill-router/docs/skill-router-overview.pptx"
    prs.save(out)
    size_kb = os.path.getsize(out) / 1024
    print(f"✅ Presentation saved: {out}")
    print(f"   File size: {size_kb:.0f} KB")
    print(f"   Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
